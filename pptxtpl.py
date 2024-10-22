#!/usr/bin/env python
#-*- coding:utf-8 -*-
import re
import copy
import slide_copy


from pptx import Presentation, parts
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from pptx.dml.color import RGBColor


class PPTXTemplate(object):
    """ Class for managing pptx files as they were self-definition templates """

    def __init__(self, pptx_template_path):
        self.presentation = Presentation(pptx_template_path)
        self.replace_label_format = "{%s}"
        self.replace_label_format_pattern = re.compile(r"\{\S+\}")
        self.slide_id_format = "{slide_id=%s}"
        self.slide_id_format_pattern = re.compile(r"\{slide_id=(\S+)\}")

    def get_replace_label_format(self):
        return self.replace_label_format

    def is_contain_replace_label(self, text):
        """
        检测text文本中是否含有 替换文本标签
        :param text:
        :return:
        """
        pattern = self.replace_label_format_pattern
        res = pattern.search(text)
        return res

    def get_replace_label_ids(self, text):
        """
        提取text文本中的replace_label_ids
        :param text:
        :return:
        """
        pattern = self.replace_label_format_pattern
        res = pattern.findall(text)
        return res

    def get_slide_single_shapes(self, index):
        """
        获取幻灯片中的所有shape，以及组合图形中的子图形
        """
        slide = self.presentation.slides[index]

        shapes = []
        for shape in slide.shapes:
            # MSO_SHAPE_TYPE.GROUP 表示组合图形
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # 获取组合图形中的元素
                for element in shape.shapes:
                    # 处理每个元素
                    shapes.append(element)
            else:
                shapes.append(shape)

        return shapes

    def add_ppt_label(self, text):
        """
        ppt中的渲染数据的名称肯定不是pptx中的标签，那么无法进行映射
        于是，对渲染数据中的名称是需要将其加上ppt标签，然后将标签写入
        ppt中，最后用self.replace_data进行替换即可
        :param text:
        :return:
        """
        return self.replace_label_format % text

    def get_replace_label_left_part(self):
        """
        获取替换标签的左边符号，比如{%s} 就获取{
        :return:
        """
        left_chs = []
        for ch in self.replace_label_format:
            if ch == "%":
                break
            left_chs.append(ch)

        return "".join(left_chs)

    def get_replace_label_right_part(self):
        """
        获取替换标签的右边符号，比如{%s} 就获取}
        :return:
        """
        right_chs = []
        flag = False
        for ch in self.replace_label_format:
            if flag:
                right_chs.append(ch)
            if ch == "s":
                flag = True
        return "".join(right_chs)

    def get_slide_run_texts(self, index):
        """
        获取指定页面的所有run_texts，如果某个标签没有在输出列表中的某个元素中，
        那么久无法进行替换
        :param index:
        :return:
        """
        slide_shapes = self.get_slide_single_shapes(index)

        run_texts = []
        for shape in slide_shapes:
            # has_text_frame 表明存在text文本，一般是文本框
            if not shape.has_text_frame:
                continue

            if not self.is_contain_replace_label(shape.text_frame.text):
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run_index, run in enumerate(paragraph.runs):
                    run_texts.append(run.text)

        print(run_texts)

    def replace_data(self, index, data):
        """
        将 data 中的数据渲染到幻灯片中
        :param index：幻灯片的顺序索引
        :param data: {"label1": "", "label2": ""}
        :return:
        """
        slide_shapes = self.get_slide_single_shapes(index)

        left_flag = self.get_replace_label_left_part()
        right_flag = self.get_replace_label_right_part()

        for shape in slide_shapes:
            # has_text_frame 表明存在text文本，一般是文本框
            if not shape.has_text_frame:
                continue

            # 如果文本框中的文本不一个标签
            if not self.is_contain_replace_label(shape.text_frame.text):
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run_index, run in enumerate(paragraph.runs):
                    # 避免出现一个替换标签 {name} 变成 起前一个text'{' 后一个text  'name}' 导致替换失败
                    if run_index - 1 >= 0 and paragraph.runs[run_index - 1].text.strip().endswith(left_flag)\
                            and not run.text.strip().startswith(left_flag):
                        paragraph.runs[run_index - 1].text = paragraph.runs[run_index - 1].text.strip().rstrip(left_flag)
                        run.text = left_flag + run.text.strip()

                    if run_index + 1 < len(paragraph.runs) and paragraph.runs[run_index + 1].text.strip().startswith(right_flag)\
                            and not run.text.strip().endswith(right_flag):
                        paragraph.runs[run_index + 1].text = paragraph.runs[run_index + 1].text.strip().lstrip(right_flag)
                        run.text = run.text.strip() + right_flag

                    replace_label_ids = self.get_replace_label_ids(run.text)
                    for replace_label_id in replace_label_ids:

                        if replace_label_id in data:
                            run.text = run.text.replace(replace_label_id, str(data[replace_label_id]))



    def get_slide_id_label_format(self):
        return self.slide_id_format

    def is_slide_id_label(self, text):
        """
        提取text文本中的slide_id_label_id
        :param text:
        :return:
        """
        pattern = self.slide_id_format_pattern
        res = pattern.match(text)
        return res

    def get_slide_id_label_id(self, text):
        """
        检测text文本是否是幻灯片id标签
        :param text:
        :return:
        """
        pattern = self.slide_id_format_pattern
        res = pattern.match(text)
        label_id = ""
        if res:
            label_id = res.group(1)
        return label_id

    def get_slide_id_index(self):
        """
        幻灯片存在多个章节，因此，可以人为用特殊标记为每个幻灯片起一个slide_id，
        当需要获取某个章节的幻灯片时就可以根据slide_id查询到幻灯片索引
        在第0页幻灯片上增加一个文本框，然后内容标记{slide_id=chapter1}
        :return:  {"chapter1": [0]}
        """
        id2indexes = {}
        for index, slide in enumerate(self.presentation):
            for shape in slide.shapes:
                # 普通的文本框
                if shape.has_text_frame:
                    if self.is_slide_id_label(shape.text_frame.text):
                        slide_id = self.get_slide_id_label_id(shape.text_frame.text)
                        id2indexes.setdefault([])
                        id2indexes[slide_id].append(index)
                        break

        return id2indexes




    @property
    def xml_slides(self):
        return list(self.presentation.slides._sldIdLst)  # pylint: disable=protected-access

    def move_slide(self, old_index, new_index):
        slides = self.xml_slides
        self.xml_slides.remove(slides[old_index])
        self.xml_slides.insert(new_index, slides[old_index])

    # also works for deleting slides
    def delete_slide(self, index):
        slides = self.xml_slides
        self.xml_slides.remove(slides[index])

    def delete_slides(self, indexes):
        # 第一页幻灯片被删除，第二页就变成第一页，所以从后往前删除更加妥当
        indexes = sorted(indexes, reverse=True)
        for index in indexes:
            self.delete_slide(index)

    def add_blank_slide(self):
        slide_layout = self.presentation.slide_layouts[0]
        self.presentation.slides.add_slide(slide_layout)

    def delete_shapes(self, index):
        """
        删除文本元素是替换标签的元素，包含文本框、统计图、组合图形
        :param index:
        :return:
        """
        slide = self.presentation.slides[index]
        for shape in slide.shapes:
            # MSO_SHAPE_TYPE.GROUP 表示组合图形
            # 删除整个组合图形
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # 获取组合图形中的元素
                for element in shape.shapes:
                    if element.has_text_frame and self.is_contain_replace_label(element.text_frame.text.strip()):
                        slide.shapes._spTree.remove(shape._element)
                        break

            # 删除普通的文本框、autoshape形状
            if shape.has_text_frame and self.is_contain_replace_label(shape.text_frame.text.strip()):
                slide.shapes._spTree.remove(shape._element)

            # 删除标题中包含替换标签的统计图
            if shape.has_chart and \
                    shape.chart.chart_title and \
                    self.is_contain_replace_label(shape.chart.chart_title.text_frame.text):
                slide.shapes._spTree.remove(shape._element)

    def delete_shapes_in_pptx(self):
        """
        删除整个幻灯片的带有替换标签的元素：文本框、统计图、组合图---暂不包括表格
        :return:
        """
        for index in range(len(self.presentation)):
            self.delete_shapes(index)





    def replace_bar_chart_data(self, index, title_data, title_replace):
        """
        幻灯片上的柱状图中的数据进行替换
        :param index: 幻灯片顺序索引
        :param title_data: 统计图标题替换标签 -> 统计图数据
        :param title_replace: 统计图标题替换标签 -> 统计图标题实际文本
        :return:
        """

        slide = self.presentation.slides[index]
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart_title = shape.chart.chart_title.text_frame.text

            if chart_title in title_data:
                self.replace_bar_chart_data_by_chart(shape.chart, list(title_data.values())[0])#title_data.get(chart_title))
                shape.chart.chart_title.text_frame.text = title_replace.get(chart_title, "")




    def replace_bar_chart_data_by_chart(self, bar_chart_obj, data):
        """
        根据的柱状图对象，填充数据，对柱状图对象进行数据替换
        :param bar_chart_obj:
        :param data: 统计图数据结构：{"category": ["不及格", "及格"], "data": {"一班": [3, 80], "二班": [2, 80]}}
        解释：上述数据表示一班和二班的及格和不及格人数的柱状图
        :return:
        """
        # define chart data ---------------------
        chart_data = CategoryChartData()
        chart_data.categories = data.get("category", [])
        for series in data.get("data", []).keys():
            chart_data.add_series(series, data.get("data", {}).get(series))

        bar_chart_obj.replace_data(chart_data)



    def add_table_data(self, index, data, font=None, font_size=Pt(18)):
        """
        填充表格数据
        注意：如果找不到对应的字体常量，可以按照策略二
        data: 二维数组，表格形式
        font: 字体：举例：'Microsoft YaHei‌'
        font_size: 字体大小
        """
        table = None
        slide = self.presentation.slides[index]
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table

        if not table:
            return
        deleted_rows = []
        for i, row in enumerate(table.rows, -1):
            if i == -1:
                continue
            if i >= len(data):
                deleted_rows.append(row)
                continue
            for j, cell in enumerate(row.cells):
                if j >= len(data[i]):
                    break

                if font:
                    # 1、策略一: 清理掉表格单元格中的文本格式及内容
                    cell.text_frame.clear()
                    p = cell.text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = str(data[i][j])
                    run.font.name = font
                    run.font.size = font_size
                else:
                # 2、策略二：如果只替换文本，但是在幻灯片中的表格的文本设置好字体及大小
                #    好处在于不用辛苦寻找字体名称
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = str(data[i][j])

        self.delete_table_rows(table, deleted_rows[::-1])

    def delete_table_rows(self, table, rows):
        """
        使用时需要注意从后往前删表格行
        """
        for row in rows:
            tbl = table._tbl
            tr = row._tr
            tbl.remove(tr)



    def get_slide_group_shapes(self, index):
        """
        获取幻灯片中的所有的组合图形
        """
        slide = self.presentation.slides[index]

        shapes = []
        for shape in slide.shapes:
            # MSO_SHAPE_TYPE.GROUP 表示组合图形
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                shapes.append(shape)

        return shapes

    def update_group_shape_position_size(self, group_shapes, position_size_list):
        """
        调整组合图形的位置和大小，需要根据方位数组进行设置
        一般移动组合图形的收益更大，但是也可以group_shapes也可以是单个图形列表
        因为组合图形和单个图形都是图形对象
        注意：group_shapes 和 position_size_list 都是一一匹配有顺序的
        :param group_shapes: 组合图形数组
        :param position_size_list: 方位数组，每个元素对应组合图形的位置和大小，例子如下：
                                    [{"left": Inches(1.14), "width": Inches(5.74)}] :return: """
        for group_shape, position_size in zip(group_shapes, position_size_list):
            if position_size.get("top"):
                group_shape.top = position_size["top"]
            if position_size.get("left"):
                group_shape.left = position_size["left"]
            if position_size.get("width"):
                group_shape.width = position_size["width"]
            if position_size.get("height"):
                group_shape.height = position_size["height"]



    def set_background_color(self, index):
        """
        设置幻灯片文本框的背景色
        根据文本框中中文本字符串的值，选择对应的背景色
        注意：文本框中已经被填充了对应的色彩字符串
        indexes: 幻灯片索引
        """

        shapes = self.get_slide_single_shapes(index)
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            # 获取文本框中的文本
            text_frame = shape.text_frame
            text = text_frame.text
            if text.strip() == 'red':
                shape.fill.solid()
                # 设置红色背景
                shape.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
                text_frame.text = ""
            elif text.strip() == 'blue':
                # 设置蓝色背景
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0xFF)
                text_frame.text = ""

    def set_text_color(self, text_shape, color):
        """
        根据执行文本框，渲染字体的颜色
        :param text_shape: 文本框图形
        :param color: 红色： RGBColor(0xFF, 0x00, 0x00)
        :return:
        """
        if text_shape.has_text_frame:
            return

        for paragraph in text_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = color



    def copy_slide(self, source_index, target_index):
        '''
        Copy ``source_index`` slide from presentation to appear at ``target_index``.
        only python-pptx=0.6.19 is supported
        chart shape is not supported
        :arg source_index: copy self.prs.presentation.slides[source_index] slide to
        :arg target_index: location to copy into. 0 makes it the first slide
        '''
        prs = self.presentation
        source = prs.slides[source_index]

        # Append slide with source's layout. Then delete shapes to get a blank slide
        dest = prs.slides.add_slide(source.slide_layout)
        for shp in dest.shapes:
            shp.element.getparent().remove(shp.element)
        # Copy shapes from source, in order
        for shape in source.shapes:
            new_shape = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
        # Copy rels from source

        for key, val in source.part.rels.items():
            target = val._target
            dest.part.rels.add_relationship(val.reltype, target, val.rId, val.is_external)
        # Move appended slide into target_index
        prs.slides.element.insert(target_index, prs.slides.element[-1])
        return dest

    def pptx_copy_slide(self, source_id, target_index):
        """
        有统计图表的复制会失败，不支持
        将 self.presentation.slides[source_id] 那一页幻灯片复制成一页新幻灯片，
        然后将新的幻灯片插入到target_index
        :param source_id:
        :param target_index:
        :return:
        """

        pres = self.presentation
        source = pres.slides[source_id]

        dest = pres.slides.add_slide(source.slide_layout)
        for shape in dest.shapes:
            shape.element.getparent().remove(shape.element)

        for shape in source.shapes:
            new_shape = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

        for key, rel in source.part.rels.items():
            target = rel._target

            if "notesSlide" in rel.reltype:
                continue

            if 'chart' in rel.reltype:
                # https://github.com/scanny/python-pptx/issues/132#issuecomment-414001942
                partname = target.package.next_partname(parts.chart.ChartPart.partname_template)
                xlsx_blob = target.chart_workbook.xlsx_part.blob
                target = parts.chart.ChartPart(
                    partname=partname,
                    content_type=target.content_type,
                    element=copy.deepcopy(target._element),
                    package=target.package)
                target.chart_workbook.xlsx_part = parts.chart.EmbeddedXlsxPart.new(
                    blob=xlsx_blob,
                    package=target.package)

            if rel.is_external:
                dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                dest.part.rels.get_or_add(rel.reltype, rel._target)

        pres.slides.element.insert(target_index, pres.slides.element[-1])

        return dest

    def duplicate_slide_with_chart(self, source_id, target_index):
        """
        复制带有统计图的幻灯片，但是目前仅仅在wps 2023上实验成功
        :param source_id:
        :param target_index:
        :return:
        """
        slide_copy.duplicate_slide(self.presentation, source_id)

        self.move_slide(source_id + 1, target_index)


    def save(self, save_path):
        self.presentation.save(save_path)

