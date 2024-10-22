#!/usr/bin/env python
#-*- coding:utf-8 -*-

from typing import Union

import pandas as pd

def chart_to_dataframe(graphical_frame) -> pd.DataFrame:
    """
    Helper to parse chart data to a DataFrame.

    :source: https://openpyxl.readthedocs.io/en/stable/pandas.html

    :param graphical_frame:
    :return:
    """
    from openpyxl import load_workbook

    from io import BytesIO
    wb = load_workbook(BytesIO(graphical_frame.chart.part.chart_workbook.xlsx_part.blob), read_only=True)

    ws = wb.active

    from itertools import islice
    import pandas as pd
    data = ws.values
    cols = next(data)[1:]
    data = list(data)
    idx = [r[0] for r in data]
    data = (islice(r, 1, None) for r in data)
    df = pd.DataFrame(data, index=idx, columns=cols)

    return df


def dataframe_to_chart_data(df):
    """
    Transforms a DataFrame to a CategoryChartData for PPT compilation.

    The indexes of the DataFrame are the categories, with each column becoming a series.

    :param df:
    :return:
    """
    from pptx.chart.data import CategoryChartData
    import numpy as np

    copy_data = CategoryChartData()
    copy_data.categories = df.index.astype(str).to_list()

    edge_cases = 0
    for c in df.columns:
        series_data = df[c].copy()
        fixed_series_data = series_data.replace([np.inf, -np.inf, np.nan], None)

        edge_cases = edge_cases + np.count_nonzero(fixed_series_data != series_data)

        copy_data.add_series(str(c), fixed_series_data.to_list())

    # Warning over data filled for compatibility
    if edge_cases > 0:
        import warnings
        warnings.warn("Series data containing NaN/INF values: filled to empty")

    return copy_data


def clone_chart(graphical_frame, dest):
    """
    Helper to clone a chart with related styling.

    :param graphical_frame:
    :param dest:
    :return:
    """
    chart = graphical_frame.chart

    df = chart_to_dataframe(graphical_frame)
    chart_data = dataframe_to_chart_data(df)

    new_chart = dest.shapes.add_chart(
        chart.chart_type,
        graphical_frame.left,
        graphical_frame.top,
        graphical_frame.width,
        graphical_frame.height,
        chart_data
    )

    # Fix offset for Graphical shape
    import copy
    cur_el = new_chart._element.xpath(".//p:nvGraphicFramePr")[0]
    ref_el = graphical_frame._element.xpath(".//p:nvGraphicFramePr")[0]
    parent = cur_el.getparent()
    parent.insert(
        parent.index(cur_el) + 1,
        copy.deepcopy(ref_el)
    )
    parent.remove(cur_el)

    # Clone styling from old chart to new one
    from random import randrange
    from lxml import etree
    from pptx.oxml import parse_xml

    id_attribute = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'

    old_chart_ref_id = graphical_frame.element.xpath(".//c:chart")[0].attrib[id_attribute]
    chart_ref_id = new_chart.element.xpath(".//c:chart")[0].attrib[id_attribute]

    new_chart_part = new_chart.part.rels._rels[chart_ref_id].target_part
    old_chart_part = graphical_frame.part.rels._rels[old_chart_ref_id].target_part

    chart_data_reference_id = new_chart_part._element.xpath(".//c:externalData")[0].attrib[id_attribute]

    cloned_styling = copy.deepcopy(old_chart_part._element)
    cloned_styling.xpath(".//c:externalData")[0].set(id_attribute, chart_data_reference_id)
    cloned_styling.xpath(".//c:autoUpdate")[0].set("val", "1")
    new_chart_part.part._element = cloned_styling

    # Parse other relationships of the chart
    from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
    from pptx.opc.package import XmlPart

    class ColorsPart(XmlPart):
        partname_template = "/ppt/charts/colors%d.xml"

        @classmethod
        def new(cls, package, element):
            part = cls.load(
                package.next_partname(cls.partname_template),
                CT.OFC_CHART_COLORS,
                package,
                element,
            )
            return part

    class StylePart(XmlPart):
        partname_template = "/ppt/charts/style%d.xml"

        @classmethod
        def new(cls, package, element):
            part = cls.load(
                package.next_partname(cls.partname_template),
                CT.OFC_CHART_STYLE,
                package,
                element,
            )
            return part

    new_chart_refs = new_chart_part.rels
    old_chart_refs = old_chart_part.rels

    # Fix styling and colors applied to the new chart
    for k, v in dict(old_chart_refs._rels).items():
        if v.reltype == 'http://schemas.microsoft.com/office/2011/relationships/chartStyle':
            targ = v.target_part

            new_el = parse_xml(copy.deepcopy(targ.blob))
            new_el.set("id", str(randrange(10 ** 5, 10 ** 9)))
            new_colors_ref = StylePart.new(targ.package, etree.tostring(new_el))
            new_chart_refs.get_or_add("http://schemas.microsoft.com/office/2011/relationships/chartStyle",
                                      new_colors_ref)
        elif v.reltype == RT.CHART_COLOR_STYLE:
            targ = v.target_part

            new_el = parse_xml(copy.deepcopy(targ.blob))
            new_el.set("id", str(randrange(10 ** 5, 10 ** 9)))
            new_colors_ref = ColorsPart.new(targ.package, etree.tostring(new_el))
            new_chart_refs.get_or_add(RT.CHART_COLOR_STYLE, new_colors_ref)

    return new_chart


def _object_rels(obj):
    rels = obj.rels

    # Change required for python-pptx 0.6.22
    check_rels_content = [k for k in rels]
    if isinstance(check_rels_content.pop(), str):
        return [v for k, v in rels.items()]
    else:
        return [k for k in rels]


def _exp_add_slide(ppt, slide_layout):
    """
    Function to handle slide creation in the Presentation, to avoid issues caused by default implementation.

    :param slide_layout:
    :return:
    """

    def generate_slide_partname(self):
        """Return |PackURI| instance containing next available slide partname."""
        from pptx.opc.packuri import PackURI

        sldIdLst = self._element.get_or_add_sldIdLst()

        existing_rels = [k.target_partname for k in _object_rels(self)]
        partname_str = "/ppt/slides/slide%d.xml" % (len(sldIdLst) + 1)

        while partname_str in existing_rels:
            import random
            import string

            random_part = ''.join(random.choice(string.ascii_letters) for i in range(2))
            partname_str = "/ppt/slides/slide%s%d.xml" % (random_part, len(sldIdLst) + 1)

        return PackURI(partname_str)

    def add_slide_part(self, slide_layout):
        """
        Return an (rId, slide) pair of a newly created blank slide that
        inherits appearance from *slide_layout*.
        """
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.parts.slide import SlidePart

        partname = generate_slide_partname(self)
        slide_layout_part = slide_layout.part
        slide_part = SlidePart.new(partname, self.package, slide_layout_part)
        rId = self.relate_to(slide_part, RT.SLIDE)
        return rId, slide_part.slide

    def add_slide_ppt(self, slide_layout):
        rId, slide = add_slide_part(self.part, slide_layout)
        slide.shapes.clone_layout_placeholders(slide_layout)
        self._sldIdLst.add_sldId(rId)
        return slide

    # slide_layout = self.get_master_slide_layout(slide_layout)
    return add_slide_ppt(ppt.slides, slide_layout)


def copy_shapes(source, dest):
    """
    Helper to copy shapes handling edge cases.

    :param source:
    :param dest:
    :return:
    """
    from pptx.shapes.group import GroupShape
    import copy

    # Copy all existing shapes
    for shape in source:
        if isinstance(shape, GroupShape):
            group = dest.shapes.add_group_shape()
            group.name = shape.name
            group.left = shape.left
            group.top = shape.top
            group.width = shape.width
            group.height = shape.height
            group.rotation = shape.rotation

            # Recursive copy of contents
            copy_shapes(shape.shapes, group)

            # Fix offset
            cur_el = group._element.xpath(".//p:grpSpPr")[0]
            ref_el = shape._element.xpath(".//p:grpSpPr")[0]
            parent = cur_el.getparent()
            parent.insert(
                parent.index(cur_el) + 1,
                copy.deepcopy(ref_el)
            )
            parent.remove(cur_el)

            result = group
        elif hasattr(shape, "image"):
            import io

            # Get image contents
            content = io.BytesIO(shape.image.blob)
            result = dest.shapes.add_picture(
                content, shape.left, shape.top, shape.width, shape.height
            )
            result.name = shape.name
            result.crop_left = shape.crop_left
            result.crop_right = shape.crop_right
            result.crop_top = shape.crop_top
            result.crop_bottom = shape.crop_bottom
        elif hasattr(shape, "has_chart") and shape.has_chart:
            result = clone_chart(shape, dest)
        else:
            import copy

            newel = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(newel, "p:extLst")
            result = dest.shapes[-1]


def duplicate_slide(ppt, slide_index: int):
    """
    Duplicate the slide with the given number in presentation.
    Adds the new slide by default at the end of the presentation.

    :param ppt:
    :param slide_index: Slide number
    :return:
    """
    source = ppt.slides[slide_index]

    dest = _exp_add_slide(ppt, source.slide_layout)

    # Remove all shapes from the default layout
    for shape in dest.shapes:
        remove_shape(shape)

    # Copy all existing shapes
    copy_shapes(source.shapes, dest)

    # Copy all existing shapes
    if source.has_notes_slide:
        txt = source.notes_slide.notes_text_frame.text
        dest.notes_slide.notes_text_frame.text = txt

    return dest

def remove_shape(shape):
    """
    Helper to remove a specific shape.

    :source: https://stackoverflow.com/questions/64700638/is-there-a-way-to-delete-a-shape-with-python-pptx

    :param shape:
    :return:
    """
    el = shape.element  # --- get reference to XML element for shape
    el.getparent().remove(el)  # --- remove that shape element from its tree

import pptx

presentation = pptx.Presentation(r"./example.pptx")
duplicate_slide(presentation, 1)

presentation.save("./save.pptx")